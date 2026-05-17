"""PowerPoint parser plugin using python-pptx.

Extracts text and images from .pptx files, organized by slide (page).
"""

import logging
import os
import uuid

from app.services.parsers.base import Asset, Block, ParsedDocument, ParseError

logger = logging.getLogger(__name__)


class PptxParser:
    """PowerPoint parser using python-pptx.

    Extracts content slide-by-slide:
    - Text from text frames (titles, body text, notes)
    - Images from shapes
    - Maintains slide number as page_number
    """

    name: str = "pptx-parser"
    supported_extensions: list[str] = ["pptx"]
    priority: int = 100

    def can_parse(self, file_path: str, mime_type: str) -> bool:
        """Check if this parser can handle the file."""
        ext = os.path.splitext(file_path)[1].lower().lstrip(".")
        return ext in self.supported_extensions or mime_type in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    async def parse(self, file_path: str) -> ParsedDocument:
        """Parse a PowerPoint file.

        Args:
            file_path: Path to the .pptx file

        Returns:
            ParsedDocument with slide-by-slide content

        Raises:
            ParseError: If the file is corrupted or password-protected
        """
        if not os.path.exists(file_path):
            raise ParseError(f"File not found: {file_path}", reason="corrupted")

        try:
            from pptx import Presentation
            from pptx.exc import PackageNotFoundError
        except ImportError:
            raise ParseError(
                "python-pptx is not installed",
                reason="unknown",
            )

        try:
            prs = Presentation(file_path)
        except PackageNotFoundError:
            raise ParseError(
                f"File is not a valid PPTX document: {file_path}",
                reason="corrupted",
            )
        except Exception as e:
            error_msg = str(e).lower()
            if "password" in error_msg or "encrypted" in error_msg:
                raise ParseError(
                    f"Presentation is password-protected: {file_path}",
                    reason="password_protected",
                )
            if "corrupt" in error_msg or "invalid" in error_msg:
                raise ParseError(
                    f"Presentation is corrupted: {file_path}",
                    reason="corrupted",
                )
            raise ParseError(
                f"Failed to parse PPTX: {file_path}: {e}",
                reason="unknown",
            )

        blocks: list[Block] = []
        assets: list[Asset] = []
        metadata: dict = {
            "source": file_path,
            "page_count": len(prs.slides),
        }

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_blocks, slide_assets = self._process_slide(slide, slide_num)
            blocks.extend(slide_blocks)
            assets.extend(slide_assets)

        if not blocks:
            raise ParseError(
                f"No content extracted from PPTX: {file_path}",
                reason="empty",
            )

        return ParsedDocument(blocks=blocks, metadata=metadata, assets=assets)

    def _process_slide(
        self, slide, slide_num: int
    ) -> tuple[list[Block], list[Asset]]:
        """Process a single slide and extract blocks and assets."""
        blocks: list[Block] = []
        assets: list[Asset] = []

        for shape in slide.shapes:
            # Handle text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue

                    # Determine block type
                    block_type = "paragraph"
                    style: dict = {}

                    # Check if this is a title placeholder
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                        ph_type = shape.placeholder_format.type
                        # Placeholder types: TITLE=1, CENTER_TITLE=3, SUBTITLE=4
                        if ph_type in (1, 3, 15):  # Title types
                            block_type = "heading"
                            style["heading_level"] = 1
                        elif ph_type == 4:  # Subtitle
                            block_type = "heading"
                            style["heading_level"] = 2

                    # Check paragraph level for indentation
                    if paragraph.level and paragraph.level > 0:
                        style["indent_level"] = paragraph.level

                    # Check for bold/italic
                    has_bold = any(
                        run.font.bold for run in paragraph.runs
                        if run.font and run.font.bold is not None
                    )
                    has_italic = any(
                        run.font.italic for run in paragraph.runs
                        if run.font and run.font.italic is not None
                    )
                    style["bold"] = has_bold
                    style["italic"] = has_italic

                    blocks.append(Block(
                        type=block_type,
                        text=text,
                        page_number=slide_num,
                        style=style,
                        raw={"shape_name": shape.name},
                    ))

            # Handle tables
            if shape.has_table:
                table_text = self._extract_table(shape.table)
                if table_text:
                    blocks.append(Block(
                        type="table",
                        text=table_text,
                        page_number=slide_num,
                        style={},
                        raw={"shape_name": shape.name},
                    ))

            # Handle images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    asset_id = str(uuid.uuid4())
                    assets.append(Asset(
                        id=asset_id,
                        type="image",
                        data=image.blob,
                        mime_type=image.content_type or "image/png",
                        page_number=slide_num,
                    ))
                    blocks.append(Block(
                        type="image",
                        text=f"[Image: {asset_id}]",
                        page_number=slide_num,
                        style={},
                        raw={"asset_id": asset_id},
                    ))
                except Exception:
                    pass

        # Extract slide notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                blocks.append(Block(
                    type="paragraph",
                    text=notes_text,
                    page_number=slide_num,
                    style={"is_notes": True},
                    raw={"source": "notes"},
                ))

        return blocks, assets

    def _extract_table(self, table) -> str:
        """Convert a PPTX table to Markdown format."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")

        if not rows:
            return ""

        # Add header separator after first row
        col_count = len(table.columns)
        separator = "| " + " | ".join(["---"] * col_count) + " |"
        rows.insert(1, separator)

        return "\n".join(rows)
